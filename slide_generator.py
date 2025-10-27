"""
AI-Powered Slide Generator
Generates PowerPoint presentations using Google Gemini AI
"""

import os
import argparse
from typing import List, Dict
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
import json

# Load environment variables
load_dotenv()


class SlideGenerator:
    """Main class for generating AI-powered presentations"""
    
    def __init__(self, api_key: str = None):
        """Initialize the slide generator with Gemini API key"""
        self.api_key = api_key or os.getenv('GEMINI_API_KEY')
        if not self.api_key:
            raise ValueError("Gemini API key not found. Please set GEMINI_API_KEY in .env file")
        
        # Use gemini-2.5-flash which is available and fast
        self.api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={self.api_key}"
    
    def generate_slide_content(self, topic: str, num_slides: int = 8, include_images: bool = False, include_code: bool = False, style: str = "professional", audience: str = "") -> List[Dict]:
        """
        Generate slide content using Google Gemini AI
        
        Args:
            topic: The presentation topic/prompt (can be detailed description)
            num_slides: Number of slides to generate
            include_images: Whether to suggest images for slides
            include_code: Whether to include code examples
            style: Presentation style (professional, educational, technical, creative, minimalist, playful)
            audience: Target audience description
            
        Returns:
            List of dictionaries containing slide information
        """
        print(f"🤖 Generating {style} presentation based on your prompt...")
        
        # Style descriptions
        style_guidance = {
            "professional": "Use formal language, focus on business value and ROI, include data and metrics",
            "educational": "Use clear explanations, include learning objectives, break down complex concepts",
            "technical": "Use technical terminology, include detailed examples, focus on implementation",
            "creative": "Use engaging language, include storytelling elements, focus on inspiration",
            "minimalist": "Use concise points, minimal text per slide, focus on key takeaways",
            "playful": "Use casual language, include fun analogies, make it engaging and light"
        }
        
        style_instruction = style_guidance.get(style, style_guidance["professional"])
        audience_instruction = f"\nTarget audience: {audience}. Tailor content complexity and language accordingly." if audience else ""
        
        # Build the prompt with proper JSON structure
        json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context"
  }'''
        
        extra_instructions = []
        if include_images:
            extra_instructions.append("- Include 'image_prompt' field with a description for relevant images (not for every slide, only where images add value)")
            json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context",
    "image_prompt": "Description of image (optional)"
  }'''
        
        if include_code:
            extra_instructions.append("- Include 'code' field with code examples when relevant to the topic (use proper syntax)")
            if include_images:
                json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context",
    "image_prompt": "Description of image (optional)",
    "code": "Code example (optional)"
  }'''
            else:
                json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context",
    "code": "Code example (optional)"
  }'''
        
        # Build the prompt with proper JSON structure
        extra_instructions = []
        if include_images:
            extra_instructions.append("- Include 'image_prompt' field with a description for relevant images (not for every slide, only where images add value)")
            json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context",
    "image_prompt": "Description of image (optional)"
  }'''
        
        if include_code:
            extra_instructions.append("- Include 'code' field with code examples when relevant to the topic (use proper syntax)")
            if include_images:
                json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context",
    "image_prompt": "Description of image (optional)",
    "code": "Code example (optional)"
  }'''
            else:
                json_example = '''  {
    "slide_number": 1,
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes with additional context",
    "code": "Code example (optional)"
  }'''
        
        prompt = f"""Create a {style} presentation with {num_slides} slides based on this user request:{audience_instruction}

USER REQUEST:
"{topic}"

STYLE GUIDANCE: {style_instruction}

For each slide, provide:
1. A clear, concise title (max 10 words)
2. 3-5 bullet points with key information (each bullet should be 1-2 sentences)
3. Speaker notes with additional context (2-3 sentences)
{chr(10).join(extra_instructions) if extra_instructions else ''}

Format your response as a JSON array with this structure:
[
{json_example}
]

Make sure:
- Slide 1 is a title slide with the presentation title and subtitle
- Content is informative, engaging, and well-structured
- Bullet points are concise but meaningful
- The presentation has a logical flow from introduction to conclusion
- The last slide is a conclusion or summary
{"- Include relevant code examples for technical concepts (use proper syntax, keep examples concise)" if include_code else ""}
{"- Suggest images that would help visualize key concepts" if include_images else ""}

Return ONLY the JSON array, no additional text."""

        try:
            # Prepare the request payload
            payload = {
                "contents": [{
                    "parts": [{
                        "text": prompt
                    }]
                }]
            }
            
            # Make the API request
            response = requests.post(self.api_url, json=payload)
            response.raise_for_status()
            
            # Parse the response
            result = response.json()
            content = result['candidates'][0]['content']['parts'][0]['text'].strip()
            
            # Extract JSON from response (in case there's extra text)
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].split("```")[0].strip()
            
            # Debug: Print the content to see what we're trying to parse
            print(f"📝 Debug - Raw JSON content (first 500 chars):\n{content[:500]}")
            
            try:
                slides_data = json.loads(content)
            except json.JSONDecodeError as je:
                print(f"❌ JSON Parse Error at line {je.lineno}, column {je.colno}")
                print(f"📝 Full content that failed to parse:\n{content}")
                raise
            
            print(f"✅ Generated {len(slides_data)} slides")
            return slides_data
            
        except Exception as e:
            print(f"❌ Error generating content: {e}")
            raise
    
    def create_presentation(self, slides_data: List[Dict], output_file: str = "presentation.pptx"):
        """
        Create a PowerPoint presentation from slide data
        
        Args:
            slides_data: List of dictionaries containing slide information
            output_file: Output filename for the presentation
        """
        print(f"📊 Creating PowerPoint presentation...")
        
        # Create presentation object
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for slide_data in slides_data:
            slide_num = slide_data.get('slide_number', 0)
            title = slide_data.get('title', 'Untitled Slide')
            bullets = slide_data.get('bullets', [])
            notes = slide_data.get('notes', '')
            code = slide_data.get('code', None)
            image_prompt = slide_data.get('image_prompt', None)
            
            if slide_num == 1:
                # Title slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                self._create_title_slide(slide, title, bullets)
            else:
                # Content slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                self._create_content_slide(slide, title, bullets, code, image_prompt)
            
            # Add speaker notes
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
        
        # Save presentation
        prs.save(output_file)
        print(f"✅ Presentation saved as '{output_file}'")
    
    def _create_title_slide(self, slide, title: str, bullets: List[str]):
        """Create a title slide with custom styling"""
        # Add background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 120)  # Dark blue
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle if available
        if bullets:
            subtitle_text = bullets[0] if len(bullets) > 0 else ""
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    def _create_content_slide(self, slide, title: str, bullets: List[str], code: str = None, image_prompt: str = None):
        """Create a content slide with title, bullet points, and optional code/images"""
        # Add background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add colored header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(31, 78, 120)
        header.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Determine layout based on content
        content_top = Inches(1.8)
        content_height = Inches(5.2)
        
        # If there's code or image, use two-column layout
        if code or image_prompt:
            # Left column for bullets
            if bullets:
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), content_top, Inches(4.5), content_height
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_before = Pt(8)
                    p.space_after = Pt(8)
            
            # Right column for code or image placeholder
            if code:
                code_box = slide.shapes.add_textbox(
                    Inches(5.2), content_top, Inches(4.3), content_height
                )
                code_frame = code_box.text_frame
                code_frame.word_wrap = True
                
                # Add code with formatting
                p = code_frame.paragraphs[0]
                p.text = code
                p.font.name = 'Courier New'
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(0, 0, 0)
                
                # Add background to code box
                code_box.fill.solid()
                code_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
                code_box.line.color.rgb = RGBColor(200, 200, 200)
            
            elif image_prompt:
                # Add image placeholder
                img_box = slide.shapes.add_textbox(
                    Inches(5.2), content_top, Inches(4.3), content_height
                )
                img_frame = img_box.text_frame
                img_frame.word_wrap = True
                img_frame.vertical_anchor = 1  # Middle
                
                p = img_frame.paragraphs[0]
                p.text = f"🎨 Image:\n{image_prompt}"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(100, 100, 100)
                p.font.italic = True
                
                # Add background
                img_box.fill.solid()
                img_box.fill.fore_color.rgb = RGBColor(245, 245, 250)
                img_box.line.color.rgb = RGBColor(200, 200, 220)
        else:
            # Full width for bullets only
            if bullets:
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), content_top, Inches(8.5), content_height
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_before = Pt(12)
                    p.space_after = Pt(12)
    
    def generate_presentation(self, topic: str, num_slides: int = 8, output_file: str = "presentation.pptx", include_images: bool = False, include_code: bool = False, style: str = "professional", audience: str = "") -> tuple:
        """
        Complete pipeline: generate content and create presentation
        
        Args:
            topic: Presentation topic
            num_slides: Number of slides to generate
            output_file: Output filename
            include_images: Whether to include image suggestions
            include_code: Whether to include code examples
            style: Presentation style
            audience: Target audience
            
        Returns:
            Tuple of (path to generated presentation, slides data)
        """
        slides_data = self.generate_slide_content(topic, num_slides, include_images, include_code, style, audience)
        self.create_presentation(slides_data, output_file)
        return output_file, slides_data


def main():
    """Command-line interface for the slide generator"""
    parser = argparse.ArgumentParser(
        description='Generate AI-powered PowerPoint presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python slide_generator.py "Introduction to Python"
  python slide_generator.py "Marketing Strategy 2024" --slides 12
  python slide_generator.py "Climate Change" --slides 10 --output climate.pptx
"""
    )
    
    parser.add_argument('topic', help='The topic of the presentation')
    parser.add_argument(
        '--slides', '-s',
        type=int,
        default=8,
        help='Number of slides to generate (default: 8)'
    )
    parser.add_argument(
        '--output', '-o',
        default='presentation.pptx',
        help='Output filename (default: presentation.pptx)'
    )
    
    args = parser.parse_args()
    
    try:
        print("=" * 60)
        print("🚀 AI Slide Generator")
        print("=" * 60)
        
        generator = SlideGenerator()
        output_path = generator.generate_presentation(
            args.topic,
            args.slides,
            args.output
        )
        
        print("=" * 60)
        print(f"🎉 Success! Presentation created: {output_path}")
        print("=" * 60)
        
    except Exception as e:
        print(f"\n❌ Error: {e}")
        return 1
    
    return 0


if __name__ == '__main__':
    exit(main())
