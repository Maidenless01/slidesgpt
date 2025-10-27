import os
import json
import requests
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
import google.generativeai as genai
import hashlib
import urllib.parse
import time

load_dotenv()

class FreeSlideGenerator:
    """Free-tier slide generator using only free APIs"""
    
    THEMES = {
        "modern_blue": {
            "name": "Modern Blue",
            "primary": RGBColor(33, 150, 243),
            "secondary": RGBColor(13, 71, 161),
            "accent": RGBColor(255, 193, 7),
            "background": RGBColor(250, 250, 250),
            "text": RGBColor(33, 33, 33),
        },
        "corporate_gray": {
            "name": "Corporate Gray",
            "primary": RGBColor(66, 66, 66),
            "secondary": RGBColor(33, 33, 33),
            "accent": RGBColor(0, 188, 212),
            "background": RGBColor(245, 245, 245),
            "text": RGBColor(66, 66, 66),
        },
        "creative_purple": {
            "name": "Creative Purple",
            "primary": RGBColor(156, 39, 176),
            "secondary": RGBColor(74, 20, 140),
            "accent": RGBColor(255, 235, 59),
            "background": RGBColor(252, 252, 252),
            "text": RGBColor(33, 33, 33),
        },
        "tech_dark": {
            "name": "Tech Dark",
            "primary": RGBColor(0, 188, 212),
            "secondary": RGBColor(0, 150, 136),
            "accent": RGBColor(255, 64, 129),
            "background": RGBColor(18, 18, 18),
            "text": RGBColor(255, 255, 255),
        },
        "elegant_gold": {
            "name": "Elegant Gold",
            "primary": RGBColor(139, 116, 61),
            "secondary": RGBColor(101, 84, 44),
            "accent": RGBColor(255, 215, 0),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(51, 51, 51),
        },
        "nature_green": {
            "name": "Nature Green",
            "primary": RGBColor(76, 175, 80),
            "secondary": RGBColor(27, 94, 32),
            "accent": RGBColor(255, 193, 7),
            "background": RGBColor(250, 250, 250),
            "text": RGBColor(33, 33, 33),
        },
        "vibrant_orange": {
            "name": "Vibrant Orange",
            "primary": RGBColor(255, 87, 34),
            "secondary": RGBColor(230, 74, 25),
            "accent": RGBColor(255, 193, 7),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(33, 33, 33),
        },
        "minimal_mono": {
            "name": "Minimal Monochrome",
            "primary": RGBColor(0, 0, 0),
            "secondary": RGBColor(97, 97, 97),
            "accent": RGBColor(189, 189, 189),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(33, 33, 33),
        },
        "sunset_gradient": {
            "name": "Sunset Gradient",
            "primary": RGBColor(255, 94, 77),
            "secondary": RGBColor(255, 145, 77),
            "accent": RGBColor(255, 209, 102),
            "background": RGBColor(255, 250, 245),
            "text": RGBColor(51, 51, 51),
        },
        "ocean_blue": {
            "name": "Ocean Blue",
            "primary": RGBColor(3, 169, 244),
            "secondary": RGBColor(1, 87, 155),
            "accent": RGBColor(0, 188, 212),
            "background": RGBColor(240, 248, 255),
            "text": RGBColor(33, 33, 33),
        }
    }
    
    def __init__(self, gemini_api_key=None):
        self.gemini_api_key = gemini_api_key or os.getenv('GEMINI_API_KEY')
        
        if not self.gemini_api_key:
            raise ValueError("GEMINI_API_KEY not found")
        
        genai.configure(api_key=self.gemini_api_key)
        self.model = genai.GenerativeModel('gemini-2.0-flash-exp')
        
        self.cache_dir = "cache"
        self.image_cache_dir = "cache/images"
        os.makedirs(self.cache_dir, exist_ok=True)
        os.makedirs(self.image_cache_dir, exist_ok=True)
    
    def get_free_image(self, search_query, width=800, height=600):
        """
        Get free images from Unsplash API (free tier: unlimited)
        """
        try:
            # Check cache first
            cache_key = hashlib.md5(f"{search_query}_{width}_{height}".encode()).hexdigest()
            cache_path = os.path.join(self.image_cache_dir, f"{cache_key}.jpg")
            
            if os.path.exists(cache_path):
                print(f"✅ Using cached image for: {search_query}")
                return cache_path
            
            print(f"🔍 Searching free image: {search_query}")
            
            # Unsplash Source (No API key needed!)
            encoded_query = urllib.parse.quote(search_query)
            url = f"https://source.unsplash.com/{width}x{height}/?{encoded_query}"
            
            response = requests.get(url, timeout=10)
            
            if response.status_code == 200:
                img = Image.open(BytesIO(response.content))
                img.save(cache_path, 'JPEG')
                print(f"✅ Downloaded and cached image")
                return cache_path
            else:
                print(f"⚠️  Image download failed, creating placeholder")
                return self.create_placeholder_image(search_query, width, height, cache_path)
                
        except Exception as e:
            print(f"⚠️  Image fetch error: {e}, creating placeholder")
            return self.create_placeholder_image(search_query, width, height, cache_path)
    
    def generate_ai_image(self, prompt, width=800, height=600):
        """
        Generate AI images using Pollinations.ai (100% FREE - no API key needed!)
        """
        try:
            # Check cache first
            cache_key = hashlib.md5(f"ai_{prompt}_{width}_{height}".encode()).hexdigest()
            cache_path = os.path.join(self.image_cache_dir, f"ai_{cache_key}.jpg")
            
            if os.path.exists(cache_path):
                print(f"✅ Using cached AI image for: {prompt}")
                return cache_path
            
            print(f"🎨 Generating AI image: {prompt}")
            
            # Pollinations.ai - Free unlimited AI image generation!
            encoded_prompt = urllib.parse.quote(prompt)
            url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true&enhance=true"
            
            response = requests.get(url, timeout=30)
            
            if response.status_code == 200:
                img = Image.open(BytesIO(response.content))
                img.save(cache_path, 'JPEG', quality=95)
                print(f"✅ AI image generated and cached")
                time.sleep(1)  # Be nice to the free API
                return cache_path
            else:
                print(f"⚠️  AI image generation failed, using stock image instead")
                return self.get_free_image(prompt, width, height)
                
        except Exception as e:
            print(f"⚠️  AI image error: {e}, using stock image instead")
            return self.get_free_image(prompt, width, height)
    
    def create_placeholder_image(self, text, width, height, save_path):
        """Create a nice placeholder image with gradient and text"""
        img = Image.new('RGB', (width, height), color=(240, 240, 245))
        draw = ImageDraw.Draw(img)
        
        # Create gradient background
        for y in range(height):
            r = int(100 + (155 * y / height))
            g = int(150 + (105 * y / height))
            b = int(200 + (55 * y / height))
            draw.line([(0, y), (width, y)], fill=(r, g, b))
        
        # Add text
        try:
            # Try to use a nice font
            font = ImageFont.truetype("arial.ttf", 40)
        except:
            font = ImageFont.load_default()
        
        # Wrap text
        words = text.split()
        lines = []
        current_line = []
        for word in words:
            current_line.append(word)
            if len(' '.join(current_line)) > 30:
                lines.append(' '.join(current_line[:-1]))
                current_line = [word]
        if current_line:
            lines.append(' '.join(current_line))
        
        # Draw text centered
        y_offset = height // 2 - (len(lines) * 25)
        for line in lines[:3]:  # Max 3 lines
            bbox = draw.textbbox((0, 0), line, font=font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            
            # Draw shadow
            draw.text((x+2, y_offset+2), line, fill=(0, 0, 0, 128), font=font)
            # Draw text
            draw.text((x, y_offset), line, fill=(255, 255, 255), font=font)
            y_offset += 50
        
        img.save(save_path, 'JPEG')
        return save_path
    
    def create_flowchart(self, slide, steps, theme_config, left=1, top=2, width=8, height=4.5):
        """Create a professional flowchart diagram"""
        # Extract text from items (handle both strings and dicts)
        def extract_text(item):
            if isinstance(item, dict):
                return item.get('text', item.get('title', item.get('name', str(item))))
            return str(item)
        
        num_steps = len(steps)
        box_height = 0.6
        spacing = (height - box_height * num_steps) / (num_steps + 1)
        
        for i, step in enumerate(steps):
            y_pos = top + spacing + i * (box_height + spacing)
            
            # Add shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(y_pos),
                Inches(width), Inches(box_height)
            )
            
            # Style the shape
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme_config['primary'] if i % 2 == 0 else theme_config['secondary']
            shape.line.color.rgb = theme_config['accent']
            shape.line.width = Pt(2)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = extract_text(step)
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add arrow to next step
            if i < num_steps - 1:
                arrow = slide.shapes.add_connector(
                    2,  # Arrow connector
                    Inches(left + width/2), Inches(y_pos + box_height),
                    Inches(left + width/2), Inches(y_pos + box_height + spacing)
                )
                arrow.line.color.rgb = theme_config['accent']
                arrow.line.width = Pt(3)
    
    def create_timeline(self, slide, events, theme_config, left=1, top=2.5, width=8):
        """Create a horizontal timeline"""
        # Extract text from items (handle both strings and dicts)
        def extract_text(item):
            if isinstance(item, dict):
                return item.get('text', item.get('title', item.get('name', str(item))))
            return str(item)
        
        num_events = len(events)
        spacing = width / (num_events - 1) if num_events > 1 else width
        
        # Draw main timeline line
        line = slide.shapes.add_connector(
            1,  # Straight line
            Inches(left), Inches(top),
            Inches(left + width), Inches(top)
        )
        line.line.color.rgb = theme_config['primary']
        line.line.width = Pt(4)
        
        for i, event in enumerate(events):
            x_pos = left + i * spacing
            
            # Add circle marker
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos - 0.15), Inches(top - 0.15),
                Inches(0.3), Inches(0.3)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = theme_config['accent']
            circle.line.color.rgb = theme_config['secondary']
            circle.line.width = Pt(2)
            
            # Add event text
            text_box = slide.shapes.add_textbox(
                Inches(x_pos - 0.6), Inches(top + 0.3),
                Inches(1.2), Inches(1)
            )
            text_frame = text_box.text_frame
            text_frame.text = extract_text(event)
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = theme_config['text']
    
    def create_comparison(self, slide, left_items, right_items, theme_config, left=0.5, top=2, width=9, height=4.5):
        """Create a two-column comparison diagram"""
        column_width = (width - 0.5) / 2
        
        # Extract text from items (handle both strings and dicts)
        def extract_text(item):
            if isinstance(item, dict):
                return item.get('text', item.get('title', item.get('name', str(item))))
            return str(item)
        
        # Left column header
        left_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(column_width), Inches(0.6)
        )
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = theme_config['primary']
        left_header.line.width = Pt(0)
        
        header_text = left_header.text_frame
        header_text.text = "Option A"
        p = header_text.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Right column header
        right_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + column_width + 0.5), Inches(top),
            Inches(column_width), Inches(0.6)
        )
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = theme_config['secondary']
        right_header.line.width = Pt(0)
        
        header_text = right_header.text_frame
        header_text.text = "Option B"
        p = header_text.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add items
        item_height = 0.5
        max_items = min(len(left_items), len(right_items), 6)
        
        for i in range(max_items):
            y_pos = top + 0.8 + i * (item_height + 0.1)
            
            # Left item
            if i < len(left_items):
                left_box = slide.shapes.add_textbox(
                    Inches(left + 0.2), Inches(y_pos),
                    Inches(column_width - 0.4), Inches(item_height)
                )
                tf = left_box.text_frame
                tf.text = "✓ " + extract_text(left_items[i])
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = theme_config['text']
            
            # Right item
            if i < len(right_items):
                right_box = slide.shapes.add_textbox(
                    Inches(left + column_width + 0.7), Inches(y_pos),
                    Inches(column_width - 0.4), Inches(item_height)
                )
                tf = right_box.text_frame
                tf.text = "✓ " + extract_text(right_items[i])
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = theme_config['text']
    
    def create_cycle_diagram(self, slide, steps, theme_config, center_x=5, center_y=4, radius=2):
        """Create a circular cycle diagram"""
        # Extract text from items (handle both strings and dicts)
        def extract_text(item):
            if isinstance(item, dict):
                return item.get('text', item.get('title', item.get('name', str(item))))
            return str(item)
        
        num_steps = len(steps)
        import math
        
        for i, step in enumerate(steps):
            # Calculate position on circle
            angle = (2 * math.pi * i / num_steps) - (math.pi / 2)  # Start at top
            x = center_x + radius * math.cos(angle)
            y = center_y + radius * math.sin(angle)
            
            # Add circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.6), Inches(y - 0.6),
                Inches(1.2), Inches(1.2)
            )
            circle.fill.solid()
            
            # Alternate colors
            colors = [theme_config['primary'], theme_config['secondary'], theme_config['accent']]
            circle.fill.fore_color.rgb = colors[i % len(colors)]
            circle.line.color.rgb = theme_config['text']
            circle.line.width = Pt(2)
            
            # Add number and text
            text_frame = circle.text_frame
            step_text = extract_text(step)
            text_frame.text = f"{i+1}\n{step_text}"
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add arrow to next step
            if num_steps > 1:
                next_angle = (2 * math.pi * ((i + 1) % num_steps) / num_steps) - (math.pi / 2)
                next_x = center_x + radius * math.cos(next_angle)
                next_y = center_y + radius * math.sin(next_angle)
                
                # Arrow between circles
                arrow_start_x = x + 0.6 * math.cos(angle)
                arrow_start_y = y + 0.6 * math.sin(angle)
                
                arrow = slide.shapes.add_connector(
                    2,  # Arrow
                    Inches(arrow_start_x), Inches(arrow_start_y),
                    Inches(next_x - 0.6 * math.cos(next_angle)), 
                    Inches(next_y - 0.6 * math.sin(next_angle))
                )
                arrow.line.color.rgb = theme_config['accent']
                arrow.line.width = Pt(3)
    
    def create_pyramid(self, slide, levels, theme_config, left=2, top=2, width=6, height=4.5):
        """Create a pyramid diagram"""
        # Extract text from items (handle both strings and dicts)
        def extract_text(item):
            if isinstance(item, dict):
                return item.get('text', item.get('title', item.get('name', str(item))))
            return str(item)
        
        num_levels = len(levels)
        
        for i, level in enumerate(levels):
            level_height = height / num_levels
            y_pos = top + i * level_height
            
            # Calculate trapezoid width (wider at bottom)
            level_width = width * (num_levels - i) / num_levels
            x_pos = left + (width - level_width) / 2
            
            # Add rectangle (approximating pyramid)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos), Inches(y_pos),
                Inches(level_width), Inches(level_height - 0.1)
            )
            
            # Style
            shape.fill.solid()
            colors = [theme_config['primary'], theme_config['secondary'], theme_config['accent']]
            shape.fill.fore_color.rgb = colors[i % len(colors)]
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.text = extract_text(level)
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
    
    def generate_slide_content(self, prompt, num_slides, style, audience, include_code, include_images, use_ai_images, theme="modern_blue"):
        """Generate slide content using Gemini (100% free)"""
        
        extra_fields = []
        if include_images:
            if use_ai_images:
                extra_fields.append('"ai_image_prompt": "detailed AI image generation prompt (e.g., professional 3D render of...)"')
            else:
                extra_fields.append('"image_search": "short search query for free stock photo"')
        if include_code:
            extra_fields.append('"code": "code example here"')
        
        # Add diagram support
        extra_fields.append('"diagram": {"type": "flowchart/timeline/comparison/cycle/pyramid", "data": [...items...]} (optional, for visual diagrams)')
        
        json_example = '{' + ','.join([
            '"title": "Slide Title"',
            '"bullets": ["Point 1", "Point 2", "Point 3"]',
            '"notes": "Speaker notes"'
        ] + extra_fields) + '}'
        
        ai_prompt = f'''Create a {num_slides}-slide presentation based on this prompt:

{prompt}

STYLE: {style}
AUDIENCE: {audience or "General audience"}
THEME: {self.THEMES[theme]["name"]}

Generate {num_slides} slides in valid JSON format. Each slide should have:
- title: Clear, engaging title
- bullets: 3-5 concise bullet points
- notes: Detailed speaker notes
{"- ai_image_prompt: DETAILED descriptive prompt for AI image generation (be specific, describe style, mood, colors)" if include_images and use_ai_images else ""}
{"- image_search: SHORT search query for free stock photos (2-4 words max)" if include_images and not use_ai_images else ""}
{"- code: Relevant code example with proper syntax" if include_code else ""}
- diagram: (optional) Visual diagram object with:
  * type: "flowchart" (for processes), "timeline" (for chronological events), "comparison" (for pros/cons), 
          "cycle" (for recurring processes), or "pyramid" (for hierarchies)
  * data: Array of items/steps for the diagram

Use diagrams strategically for:
- Flowchart: Step-by-step processes, workflows
- Timeline: Historical events, project phases, roadmaps
- Comparison: Pros vs cons, before vs after, two options
- Cycle: Iterative processes, life cycles, continuous improvement
- Pyramid: Hierarchies, priorities, levels

Return ONLY a JSON object like:
{{
  "slides": [
    {json_example},
    {json_example}
  ]
}}

Make the content engaging, well-structured, and visually rich with diagrams where appropriate.
{"For AI images, be very descriptive: 'professional 3D render of a modern office workspace, bright natural lighting, minimalist design, blue and white colors'" if use_ai_images else ""}
Include at least 2-3 diagrams throughout the presentation for better visual engagement.'''

        try:
            print(f"🤖 Generating content with Gemini ({theme} theme)...")
            response = self.model.generate_content(ai_prompt)
            
            # Clean and parse JSON
            content = response.text.strip()
            if content.startswith('```json'):
                content = content[7:]
            if content.startswith('```'):
                content = content[3:]
            if content.endswith('```'):
                content = content[:-3]
            content = content.strip()
            
            data = json.loads(content)
            
            if "slides" not in data:
                raise ValueError("Response missing 'slides' key")
            
            print(f"✅ Generated {len(data['slides'])} slides")
            return data["slides"]
            
        except json.JSONDecodeError as e:
            print(f"❌ JSON parsing error: {e}")
            print(f"Response content: {content[:500]}")
            raise
        except Exception as e:
            print(f"❌ Error: {e}")
            raise
    
    def create_presentation(self, slides_data, output_path, theme="modern_blue", use_images=False, use_ai_images=False):
        """Create PowerPoint with diagrams, free images, and AI-generated images"""
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        theme_config = self.THEMES[theme]
        
        print(f"🎨 Creating presentation with {theme_config['name']} theme...")
        
        for idx, slide_data in enumerate(slides_data):
            print(f"📄 Processing slide {idx + 1}/{len(slides_data)}: {slide_data.get('title', 'Untitled')}")
            
            # Create blank slide
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Add background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme_config['background']
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get('title', 'Untitled')
            
            title_p = title_frame.paragraphs[0]
            title_p.font.size = Pt(44)
            title_p.font.bold = True
            title_p.font.color.rgb = theme_config['primary']
            
            # Decorative line
            line = slide.shapes.add_shape(
                1,  # Line shape
                Inches(0.5), Inches(1.4), Inches(9), Inches(0)
            )
            line.line.color.rgb = theme_config['accent']
            line.line.width = Pt(3)
            
            # Check for diagram
            has_diagram = 'diagram' in slide_data and slide_data['diagram']
            
            # Get image if requested
            image_path = None
            if use_images:
                if use_ai_images and 'ai_image_prompt' in slide_data and slide_data['ai_image_prompt']:
                    image_path = self.generate_ai_image(slide_data['ai_image_prompt'], 1024, 768)
                elif 'image_search' in slide_data and slide_data['image_search']:
                    image_path = self.get_free_image(slide_data['image_search'], 800, 600)
            
            # Layout logic
            content_top = 1.7
            content_left = 0.5
            content_width = 9
            has_image = image_path and os.path.exists(image_path)
            
            if has_diagram:
                # Diagram takes center stage
                diagram_data = slide_data['diagram']
                diagram_type = diagram_data.get('type', 'flowchart')
                diagram_items = diagram_data.get('data', [])
                
                print(f"  📊 Adding {diagram_type} diagram with {len(diagram_items)} items")
                
                # Add diagram based on type
                if diagram_type == 'flowchart' and diagram_items:
                    self.create_flowchart(slide, diagram_items, theme_config)
                elif diagram_type == 'timeline' and diagram_items:
                    self.create_timeline(slide, diagram_items, theme_config)
                elif diagram_type == 'comparison' and len(diagram_items) >= 2:
                    mid = len(diagram_items) // 2
                    self.create_comparison(slide, diagram_items[:mid], diagram_items[mid:], theme_config)
                elif diagram_type == 'cycle' and diagram_items:
                    self.create_cycle_diagram(slide, diagram_items, theme_config)
                elif diagram_type == 'pyramid' and diagram_items:
                    self.create_pyramid(slide, diagram_items, theme_config)
                
                # Add bullets below or skip if diagram is main content
                if slide_data.get('bullets') and len(slide_data['bullets']) <= 3:
                    text_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(6.2),
                        Inches(9), Inches(1)
                    )
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    for bullet in slide_data['bullets'][:3]:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(14)
                        p.font.color.rgb = theme_config['text']
                
            elif has_image:
                # Two-column layout with image
                content_width = 5.2
                
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(5.5), Inches(content_top),
                        width=Inches(4), height=Inches(4.5)
                    )
                    print(f"  🖼️  Added image: {os.path.basename(image_path)}")
                except Exception as e:
                    print(f"  ⚠️  Could not add image: {e}")
                    content_width = 9
                
                # Content bullets
                if slide_data.get('bullets'):
                    text_box = slide.shapes.add_textbox(
                        Inches(content_left), Inches(content_top),
                        Inches(content_width), Inches(4.5)
                    )
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    for bullet in slide_data['bullets']:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.color.rgb = theme_config['text']
                        p.space_before = Pt(12)
            
            else:
                # Full-width content
                if slide_data.get('bullets'):
                    text_box = slide.shapes.add_textbox(
                        Inches(content_left), Inches(content_top),
                        Inches(content_width), Inches(4.5)
                    )
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    for bullet in slide_data['bullets']:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(22)
                        p.font.color.rgb = theme_config['text']
                        p.space_before = Pt(14)
            
            # Code block (if no diagram)
            if not has_diagram and 'code' in slide_data and slide_data['code']:
                code_top = 5.8 if has_image else 5.5
                code_box = slide.shapes.add_textbox(
                    Inches(content_left), Inches(code_top),
                    Inches(content_width), Inches(1.5)
                )
                code_frame = code_box.text_frame
                code_frame.text = slide_data['code']
                
                code_p = code_frame.paragraphs[0]
                code_p.font.size = Pt(12)
                code_p.font.name = "Courier New"
                code_p.font.color.rgb = RGBColor(248, 248, 242)
                
                # Code background
                code_fill = code_box.fill
                code_fill.solid()
                code_fill.fore_color.rgb = RGBColor(40, 42, 54)
            
            # Speaker notes
            if slide_data.get('notes'):
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.text = slide_data['notes']
        
        prs.save(output_path)
        print(f"✅ Saved presentation: {output_path}")
    
    def generate_presentation(self, prompt, num_slides, output_path, **options):
        """Generate complete presentation - 100% FREE with diagrams and AI images"""
        
        style = options.get('style', 'professional')
        audience = options.get('audience', '')
        include_code = options.get('include_code', False)
        include_images = options.get('include_images', False)
        use_ai_images = options.get('use_ai_images', False)  # NEW: AI-generated images
        theme = options.get('theme', 'modern_blue')
        
        # Generate content (FREE - Gemini)
        slides_data = self.generate_slide_content(
            prompt, num_slides, style, audience, include_code, include_images, use_ai_images, theme
        )
        
        # Save slides data to cache
        cache_filename = os.path.basename(output_path).replace('.pptx', '.json')
        cache_path = os.path.join(self.cache_dir, cache_filename)
        with open(cache_path, 'w', encoding='utf-8') as f:
            json.dump(slides_data, f, indent=2, ensure_ascii=False)
        
        # Create presentation with FREE images and diagrams
        self.create_presentation(slides_data, output_path, theme, include_images, use_ai_images)
        
        return {
            'success': True,
            'output_path': output_path,
            'slides_data': slides_data,
            'theme': theme,
            'num_slides': len(slides_data)
        }


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Free AI Slide Generator with Diagrams & AI Images')
    parser.add_argument('prompt', help='Presentation prompt/topic')
    parser.add_argument('--slides', type=int, default=8, help='Number of slides')
    parser.add_argument('--output', default='presentation.pptx', help='Output file')
    parser.add_argument('--theme', default='modern_blue', choices=list(FreeSlideGenerator.THEMES.keys()))
    parser.add_argument('--style', default='professional')
    parser.add_argument('--audience', default='')
    parser.add_argument('--code', action='store_true', help='Include code examples')
    parser.add_argument('--images', action='store_true', help='Include stock images')
    parser.add_argument('--ai-images', action='store_true', help='Use AI-generated images instead of stock photos')
    
    args = parser.parse_args()
    
    generator = FreeSlideGenerator()
    
    result = generator.generate_presentation(
        args.prompt,
        args.slides,
        args.output,
        style=args.style,
        audience=args.audience,
        include_code=args.code,
        include_images=args.images or args.ai_images,
        use_ai_images=args.ai_images,
        theme=args.theme
    )
    
    print(f"\n🎉 Success! Created {result['num_slides']} slides with {result['theme']} theme")
    print(f"📁 Saved to: {result['output_path']}")
    print(f"✨ Features: Diagrams, {'AI Images' if args.ai_images else 'Stock Photos' if args.images else 'No Images'}")

